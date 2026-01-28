"""
리포트 생성기 (Report Generator)
================================
AMOREPACIFIC IR 스타일의 고품질 인사이트 리포트 생성

## 지원 포맷
- DOCX: 편집 가능한 Word 문서 (기본)
- PDF: 공유용 읽기 전용 (WeasyPrint)
- PPTX: 발표용 슬라이드

## 디자인 시스템
- Pacific Blue: #001C58 (헤더, 제목)
- Amore Blue: #1F5795 (강조, 링크)
- Gray: #7D7D7D (보조 텍스트)
- White: #FFFFFF (배경)
- Accent Red: #E53935 (하락, 경고)
- Accent Green: #43A047 (상승, 긍정)

## 레이아웃
- 표지 페이지 (별도)
- 목차 페이지 (별도)
- 본문 섹션
- 참고자료
"""

import logging
import tempfile
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any

from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import parse_xml
from docx.oxml.ns import nsdecls
from docx.shared import Cm, Inches, Pt, RGBColor

logger = logging.getLogger(__name__)


# =============================================================================
# Design System Constants
# =============================================================================


class DesignSystem:
    """AMOREPACIFIC 디자인 시스템"""

    # Colors
    PACIFIC_BLUE = RGBColor(0, 28, 88)  # #001C58
    AMORE_BLUE = RGBColor(31, 87, 149)  # #1F5795
    GRAY = RGBColor(125, 125, 125)  # #7D7D7D
    LIGHT_GRAY = RGBColor(245, 245, 245)  # #F5F5F5
    WHITE = RGBColor(255, 255, 255)  # #FFFFFF
    ACCENT_RED = RGBColor(229, 57, 53)  # #E53935
    ACCENT_GREEN = RGBColor(67, 160, 71)  # #43A047

    # Font Family - Arita Dotum (돋움, 산세리프 - 제목/목차용)
    FONT_DOTUM = "Arita Dotum KR"
    FONT_DOTUM_MEDIUM = "Arita Dotum KR Medium"
    FONT_DOTUM_SEMIBOLD = "Arita Dotum KR SemiBold"
    FONT_DOTUM_BOLD = "Arita Dotum KR Bold"
    FONT_DOTUM_LIGHT = "Arita Dotum KR Light"

    # Font Family - Arita Buri (부리, 세리프 - 본문용)
    FONT_BURI = "Arita Buri KR"
    FONT_BURI_MEDIUM = "Arita Buri KR Medium"
    FONT_BURI_SEMIBOLD = "Arita Buri KR SemiBold"
    FONT_BURI_BOLD = "Arita Buri KR Bold"
    FONT_BURI_LIGHT = "Arita Buri KR Light"

    # 레거시 호환성
    FONT_FAMILY = FONT_DOTUM
    FONT_FAMILY_MEDIUM = FONT_DOTUM_MEDIUM
    FONT_FAMILY_SEMIBOLD = FONT_DOTUM_SEMIBOLD
    FONT_FAMILY_BOLD = FONT_DOTUM_BOLD
    FONT_FAMILY_LIGHT = FONT_DOTUM_LIGHT
    FONT_FAMILY_FALLBACK = "Malgun Gothic"  # Windows fallback
    FONT_FAMILY_MAC = "Apple SD Gothic Neo"  # macOS fallback

    # Typography (points)
    TITLE_SIZE = Pt(28)
    SUBTITLE_SIZE = Pt(16)
    HEADING1_SIZE = Pt(18)
    HEADING2_SIZE = Pt(14)
    HEADING3_SIZE = Pt(12)
    BODY_SIZE = Pt(11)
    CAPTION_SIZE = Pt(9)

    # Spacing
    PAGE_MARGIN_TOP = Cm(2)
    PAGE_MARGIN_BOTTOM = Cm(2)
    PAGE_MARGIN_LEFT = Cm(2.5)
    PAGE_MARGIN_RIGHT = Cm(2.5)

    # Logo paths (relative to project root)
    LOGO_COLOR_EN = "Amorepacific_CI_wordmark_v1.00/png/Color/Amorepacific_Wordmark_Color_En.png"
    LOGO_BASIC_EN = "Amorepacific_CI_wordmark_v1.00/png/Basic/Amorepacific_Wordmark_Basic_En.png"
    LOGO_REVERSE_EN = (
        "Amorepacific_CI_wordmark_v1.00/png/Reverse/Amorepacific_Wordmark_Reverse_En.png"
    )


# =============================================================================
# DOCX Report Generator
# =============================================================================


class DocxReportGenerator:
    """
    DOCX 리포트 생성기

    AMOREPACIFIC IR 스타일의 고품질 Word 문서 생성
    """

    def __init__(self, output_dir: str = None):
        """
        Args:
            output_dir: 출력 디렉토리 (None이면 임시 폴더)
        """
        self.output_dir = Path(output_dir) if output_dir else Path(tempfile.gettempdir())
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Project root 찾기
        self.project_root = self._find_project_root()

    def _find_project_root(self) -> Path:
        """프로젝트 루트 디렉토리 찾기"""
        current = Path(__file__).resolve()
        for parent in current.parents:
            if (parent / "CLAUDE.md").exists() or (parent / "dashboard_api.py").exists():
                return parent
        return Path.cwd()

    def _get_logo_path(self, logo_type: str = "color") -> Path | None:
        """로고 파일 경로 반환"""
        logo_map = {
            "color": DesignSystem.LOGO_COLOR_EN,
            "basic": DesignSystem.LOGO_BASIC_EN,
            "reverse": DesignSystem.LOGO_REVERSE_EN,
        }
        logo_rel = logo_map.get(logo_type, DesignSystem.LOGO_COLOR_EN)
        logo_path = self.project_root / logo_rel

        if logo_path.exists():
            return logo_path

        logger.warning(f"Logo not found: {logo_path}")
        return None

    def _setup_document_styles(self, doc: Document) -> None:
        """문서 기본 스타일 설정 (Arita 폰트 적용)"""
        # Normal 스타일 - Arita Buri KR Medium (본문용 - 세리프)
        style = doc.styles["Normal"]
        font = style.font
        font.name = DesignSystem.FONT_BURI_MEDIUM
        font.size = DesignSystem.BODY_SIZE

        # 단락 스타일
        paragraph_format = style.paragraph_format
        paragraph_format.space_after = Pt(6)
        paragraph_format.line_spacing = 1.15

        # Heading 1 스타일 - Arita Dotum KR Bold (제목용 - 산세리프)
        if "Heading 1" in doc.styles:
            h1_style = doc.styles["Heading 1"]
            h1_style.font.name = DesignSystem.FONT_DOTUM_BOLD
            h1_style.font.size = DesignSystem.HEADING1_SIZE
            h1_style.font.color.rgb = DesignSystem.PACIFIC_BLUE

        # Heading 2 스타일 - Arita Dotum KR SemiBold (소제목용 - 산세리프)
        if "Heading 2" in doc.styles:
            h2_style = doc.styles["Heading 2"]
            h2_style.font.name = DesignSystem.FONT_DOTUM_SEMIBOLD
            h2_style.font.size = DesignSystem.HEADING2_SIZE
            h2_style.font.color.rgb = DesignSystem.AMORE_BLUE

    def _setup_page_margins(self, doc: Document) -> None:
        """페이지 여백 설정"""
        for section in doc.sections:
            section.top_margin = DesignSystem.PAGE_MARGIN_TOP
            section.bottom_margin = DesignSystem.PAGE_MARGIN_BOTTOM
            section.left_margin = DesignSystem.PAGE_MARGIN_LEFT
            section.right_margin = DesignSystem.PAGE_MARGIN_RIGHT

    def _add_header_bar(self, doc: Document) -> None:
        """
        상단 Pacific Blue 헤더 바 추가 (테이블 활용)

        Note: python-docx는 도형(shape) 직접 지원이 제한적이므로
        색상이 있는 테이블로 헤더 바 효과 구현
        """
        # 전체 너비 테이블 (1행 1열)
        table = doc.add_table(rows=1, cols=1)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER

        # 테이블 너비를 페이지 전체로
        table.autofit = False
        table.allow_autofit = False

        cell = table.rows[0].cells[0]
        cell.width = Inches(7.5)  # A4 기준 약 페이지 너비

        # 셀 배경색 설정
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="001C58"/>')
        tcPr.append(shading)

        # 셀 높이 설정 (얇은 바)
        cell.paragraphs[0].clear()
        cell.paragraphs[0].add_run(" ")  # 최소 높이 확보

        # 테이블 행 높이 설정
        tr = table.rows[0]._tr
        trPr = tr.get_or_add_trPr()
        trHeight = parse_xml(f'<w:trHeight {nsdecls("w")} w:val="200" w:hRule="exact"/>')
        trPr.append(trHeight)

    def _add_cover_page(
        self,
        doc: Document,
        title: str,
        subtitle: str = None,
        date_range: str = None,
        generation_date: str = None,
    ) -> None:
        """
        표지 페이지 추가

        Args:
            doc: Document 객체
            title: 메인 제목
            subtitle: 부제목 (선택)
            date_range: 분석 기간 (예: "2026-01-01 ~ 2026-01-14")
            generation_date: 생성 일시
        """
        # 상단 헤더 바
        self._add_header_bar(doc)

        # 여백
        doc.add_paragraph()
        doc.add_paragraph()
        doc.add_paragraph()

        # 로고 (TOC 페이지와 동일한 크기)
        logo_path = self._get_logo_path("color")
        if logo_path:
            logo_para = doc.add_paragraph()
            logo_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = logo_para.add_run()
            run.add_picture(str(logo_path), width=Cm(10))

        # 큰 여백
        for _ in range(4):
            doc.add_paragraph()

        # 메인 제목 (두 줄로 분리: "LANEIGE Amazon US" + "경쟁력 분석 보고서")
        # 제목에 "경쟁력 분석 보고서" 또는 "분석 보고서"가 포함되어 있으면 분리
        title_line1 = title
        title_line2 = None

        # 한국어 제목 분리 패턴
        split_keywords = ["경쟁력 분석 보고서", "분석 보고서", "보고서"]
        for keyword in split_keywords:
            if keyword in title:
                parts = title.split(keyword, 1)
                title_line1 = parts[0].strip()
                title_line2 = keyword
                break

        # 첫 번째 줄
        title_para1 = doc.add_paragraph()
        title_para1.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_run1 = title_para1.add_run(title_line1)
        title_run1.font.size = DesignSystem.TITLE_SIZE
        title_run1.font.bold = True
        title_run1.font.color.rgb = DesignSystem.PACIFIC_BLUE
        title_run1.font.name = DesignSystem.FONT_FAMILY_BOLD

        # 두 번째 줄 (있는 경우)
        if title_line2:
            title_para2 = doc.add_paragraph()
            title_para2.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_run2 = title_para2.add_run(title_line2)
            title_run2.font.size = DesignSystem.TITLE_SIZE
            title_run2.font.bold = True
            title_run2.font.color.rgb = DesignSystem.PACIFIC_BLUE
            title_run2.font.name = DesignSystem.FONT_FAMILY_BOLD

        # 부제목
        if subtitle:
            subtitle_para = doc.add_paragraph()
            subtitle_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            subtitle_run = subtitle_para.add_run(subtitle)
            subtitle_run.font.size = DesignSystem.SUBTITLE_SIZE
            subtitle_run.font.color.rgb = DesignSystem.AMORE_BLUE

        # 여백
        doc.add_paragraph()
        doc.add_paragraph()

        # 분석 기간
        if date_range:
            date_para = doc.add_paragraph()
            date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            date_run = date_para.add_run(f"분석 기간: {date_range}")
            date_run.font.size = DesignSystem.BODY_SIZE
            date_run.font.color.rgb = DesignSystem.GRAY

        # 생성 일시
        if generation_date:
            gen_para = doc.add_paragraph()
            gen_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            gen_run = gen_para.add_run(f"생성일시: {generation_date}")
            gen_run.font.size = DesignSystem.BODY_SIZE
            gen_run.font.color.rgb = DesignSystem.GRAY

        # 페이지 브레이크
        doc.add_page_break()

    def _add_toc_page(
        self,
        doc: Document,
        sections: list[str],
        disclaimer: str = None,
    ) -> None:
        """
        목차 페이지 추가

        Args:
            doc: Document 객체
            sections: 섹션 목록 ["1. Executive Summary", "2. LANEIGE 분석", ...]
            disclaimer: 면책 조항 (선택)
        """
        # 상단 헤더 바
        self._add_header_bar(doc)

        # 여백
        doc.add_paragraph()

        # 로고 이미지 (AMORE PACIFIC) - 중앙 정렬, 크게
        logo_path = self._get_logo_path("color")
        if logo_path:
            logo_para = doc.add_paragraph()
            logo_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = logo_para.add_run()
            run.add_picture(str(logo_path), width=Cm(10))  # 로고 크기 증가

        # 로고와 CONTENTS 사이 여백 증가
        for _ in range(3):
            doc.add_paragraph()

        # CONTENTS 제목
        contents_para = doc.add_paragraph()
        contents_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        contents_run = contents_para.add_run("CONTENTS")
        contents_run.font.size = Pt(24)
        contents_run.font.bold = True
        contents_run.font.color.rgb = DesignSystem.PACIFIC_BLUE
        contents_run.font.name = "Arita Dotum KR SemiBold"  # Arita 폰트

        # CONTENTS와 목차 사이 여백
        spacer = doc.add_paragraph()
        spacer.paragraph_format.space_after = Pt(24)

        # 목차 항목들 (번호 + 섹션명, 왼쪽 정렬)
        for i, section in enumerate(sections, 1):
            toc_para = doc.add_paragraph()
            toc_para.alignment = WD_ALIGN_PARAGRAPH.LEFT
            toc_para.paragraph_format.left_indent = Cm(3.0)  # 왼쪽 들여쓰기
            toc_para.paragraph_format.space_before = Pt(8)
            toc_para.paragraph_format.space_after = Pt(8)

            # 번호
            num_run = toc_para.add_run(f"{i}  ")
            num_run.font.size = Pt(12)
            num_run.font.bold = True
            num_run.font.color.rgb = DesignSystem.PACIFIC_BLUE
            num_run.font.name = "Arita Dotum KR SemiBold"

            # 섹션명 (번호 제외)
            section_name = section.split(". ", 1)[-1] if ". " in section else section
            name_run = toc_para.add_run(section_name)
            name_run.font.size = Pt(12)
            name_run.font.color.rgb = DesignSystem.PACIFIC_BLUE
            name_run.font.name = "Arita Dotum KR Medium"

        # DISCLAIMER (페이지 최하단에 위치하도록 여백 추가)
        if disclaimer:
            # 충분한 여백으로 하단 배치 (목차 항목 수에 따라 조정)
            remaining_space = max(2, 10 - len(sections))
            for _ in range(remaining_space):
                doc.add_paragraph()

            disc_heading = doc.add_paragraph()
            disc_heading.paragraph_format.left_indent = Cm(1.0)
            disc_run = disc_heading.add_run("DISCLAIMER")
            disc_run.font.size = Pt(9)
            disc_run.font.bold = True
            disc_run.font.color.rgb = DesignSystem.GRAY

            disc_para = doc.add_paragraph()
            disc_para.paragraph_format.space_before = Pt(4)
            disc_para.paragraph_format.left_indent = Cm(1.0)
            disc_text = disc_para.add_run(disclaimer)
            disc_text.font.size = Pt(8)
            disc_text.font.color.rgb = DesignSystem.GRAY
            disc_text.font.name = "Arita Dotum KR Light"

        # 페이지 브레이크
        doc.add_page_break()

    def _add_section_heading(
        self,
        doc: Document,
        section_id: int,
        title: str,
        level: int = 1,
    ) -> None:
        """
        섹션 제목 추가

        Args:
            doc: Document 객체
            section_id: 섹션 번호
            title: 섹션 제목
            level: 제목 레벨 (1, 2, 3)
        """
        if level == 1:
            heading = doc.add_heading(f"{section_id}. {title}", level=1)
            for run in heading.runs:
                run.font.color.rgb = DesignSystem.PACIFIC_BLUE
                run.font.size = DesignSystem.HEADING1_SIZE
        elif level == 2:
            heading = doc.add_heading(title, level=2)
            for run in heading.runs:
                run.font.color.rgb = DesignSystem.AMORE_BLUE
                run.font.size = DesignSystem.HEADING2_SIZE
        else:
            heading = doc.add_heading(title, level=3)
            for run in heading.runs:
                run.font.color.rgb = DesignSystem.GRAY
                run.font.size = DesignSystem.HEADING3_SIZE

    def _add_content_paragraph(
        self,
        doc: Document,
        text: str,
        is_highlight: bool = False,
        is_bullet: bool = False,
        add_space_before: bool = False,
    ) -> None:
        """
        본문 단락 추가

        Args:
            doc: Document 객체
            text: 텍스트 내용
            is_highlight: 강조 여부 (■로 시작하는 항목 - 목차)
            is_bullet: 불릿 포인트 여부
            add_space_before: 앞에 빈 줄 추가 여부
        """
        # 목차(■) 앞에 빈 줄 추가
        if add_space_before:
            doc.add_paragraph()

        if is_bullet:
            para = doc.add_paragraph(style="List Bullet")
            run = para.add_run(text)
        else:
            para = doc.add_paragraph()
            run = para.add_run(text)

        if is_highlight:
            # 목차(■)는 아리따 돋움 (산세리프)
            run.font.bold = True
            run.font.color.rgb = DesignSystem.PACIFIC_BLUE
            run.font.name = DesignSystem.FONT_DOTUM_SEMIBOLD
            run.font.size = DesignSystem.HEADING3_SIZE
        else:
            # 본문은 아리따 부리 (세리프)
            run.font.name = DesignSystem.FONT_BURI_MEDIUM
            run.font.size = DesignSystem.BODY_SIZE

    def _add_kpi_cards(
        self,
        doc: Document,
        kpis: list[dict[str, Any]],
    ) -> None:
        """
        KPI 카드 추가 (테이블 기반)

        Args:
            doc: Document 객체
            kpis: [{"name": "SoS", "value": "2.5%", "change": "+0.3%", "trend": "up"}, ...]
        """
        if not kpis:
            return

        # 3열 테이블
        num_cols = min(3, len(kpis))
        table = doc.add_table(rows=1, cols=num_cols)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER

        for i, kpi in enumerate(kpis[:3]):
            cell = table.rows[0].cells[i]

            # 셀 스타일링 (배경색)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="F5F5F5"/>')
            tcPr.append(shading)

            # KPI 이름
            name_para = cell.paragraphs[0]
            name_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            name_run = name_para.add_run(kpi.get("name", ""))
            name_run.font.size = Pt(10)
            name_run.font.color.rgb = DesignSystem.GRAY

            # KPI 값
            value_para = cell.add_paragraph()
            value_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            value_run = value_para.add_run(str(kpi.get("value", "-")))
            value_run.font.size = Pt(24)
            value_run.font.bold = True
            value_run.font.color.rgb = DesignSystem.PACIFIC_BLUE

            # 변화율
            if kpi.get("change"):
                change_para = cell.add_paragraph()
                change_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                change_text = kpi["change"]
                change_run = change_para.add_run(change_text)
                change_run.font.size = Pt(12)

                # 색상 (상승/하락)
                if kpi.get("trend") == "up" or change_text.startswith("+"):
                    change_run.font.color.rgb = DesignSystem.ACCENT_GREEN
                elif kpi.get("trend") == "down" or change_text.startswith("-"):
                    change_run.font.color.rgb = DesignSystem.ACCENT_RED
                else:
                    change_run.font.color.rgb = DesignSystem.GRAY

        doc.add_paragraph()

    def _add_chart_image(
        self,
        doc: Document,
        chart_path: Path | str,
        caption: str = None,
        width: float = 6.0,
    ) -> None:
        """
        차트 이미지 추가

        Args:
            doc: Document 객체
            chart_path: 차트 이미지 파일 경로
            caption: 캡션 (선택)
            width: 이미지 너비 (인치)
        """
        chart_path = Path(chart_path)
        if not chart_path.exists():
            logger.warning(f"Chart not found: {chart_path}")
            return

        # 이미지 추가
        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = para.add_run()
        run.add_picture(str(chart_path), width=Inches(width))

        # 캡션
        if caption:
            caption_para = doc.add_paragraph()
            caption_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            caption_run = caption_para.add_run(caption)
            caption_run.font.size = DesignSystem.CAPTION_SIZE
            caption_run.font.color.rgb = DesignSystem.GRAY
            caption_run.font.italic = True

    def _add_footer(self, doc: Document) -> None:
        """푸터 추가"""
        doc.add_paragraph()
        footer = doc.add_paragraph()
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = footer.add_run(f"© {datetime.now().year} AMORE Pacific - Confidential")
        run.font.size = DesignSystem.CAPTION_SIZE
        run.font.color.rgb = DesignSystem.GRAY
        run.font.italic = True

    def generate_analyst_report(
        self,
        report_data: dict[str, Any],
        chart_paths: dict[str, Path] = None,
        output_filename: str = None,
    ) -> Path:
        """
        애널리스트 리포트 생성

        Args:
            report_data: 리포트 데이터
                {
                    "title": "LANEIGE Amazon US 경쟁력 분석 보고서",
                    "subtitle": "Weekly Insight Report",
                    "start_date": "2026-01-01",
                    "end_date": "2026-01-14",
                    "sections": [
                        {
                            "id": 1,
                            "title": "Executive Summary",
                            "content": "...",
                            "kpis": [{"name": "SoS", "value": "2.5%", ...}]
                        },
                        ...
                    ],
                    "references": "..."
                }
            chart_paths: 차트 이미지 경로 {"sos_trend": Path, ...}
            output_filename: 출력 파일명 (None이면 자동 생성)

        Returns:
            생성된 파일 경로
        """
        doc = Document()

        # 기본 설정
        self._setup_document_styles(doc)
        self._setup_page_margins(doc)

        # 메타 정보
        title = report_data.get("title", "AMORE Insight Report")
        subtitle = report_data.get("subtitle")
        start_date = report_data.get("start_date", "")
        end_date = report_data.get("end_date", "")
        date_range = f"{start_date} ~ {end_date}" if start_date and end_date else None
        generation_date = datetime.now().strftime("%Y-%m-%d %H:%M")

        # 1. 표지 페이지
        self._add_cover_page(
            doc,
            title=title,
            subtitle=subtitle,
            date_range=date_range,
            generation_date=generation_date,
        )

        # 2. 목차 페이지
        sections = report_data.get("sections", [])
        section_titles = [
            f"{s.get('id', i+1)}. {s.get('title', '')}" for i, s in enumerate(sections)
        ]

        disclaimer = (
            "본 자료는 Amazon.com 웹사이트에서 공개적으로 수집된 베스트셀러 순위 데이터를 기반으로 작성되었습니다. "
            "모든 분석 인사이트는 AI(인공지능)에 의해 자동 생성된 것으로, 실제 시장 상황과 다를 수 있습니다. "
            "본 자료의 정보를 활용한 의사결정에 대해 당사는 어떠한 책임도 지지 않습니다."
        )

        self._add_toc_page(doc, section_titles, disclaimer=disclaimer)

        # 3. 본문 섹션들
        for section in sections:
            section_id = section.get("id", 1)
            section_title = section.get("title", "")
            content = section.get("content", "")
            kpis = section.get("kpis", [])
            chart_key = section.get("chart_key")  # "sos_trend", "brand_comparison" 등

            # 섹션 제목
            self._add_section_heading(doc, section_id, section_title)

            # KPI 카드 (있는 경우)
            if kpis:
                self._add_kpi_cards(doc, kpis)

            # 본문 내용
            if content:
                is_first_heading = True  # 첫 번째 목차(■)는 줄바꿈 안함
                for line in content.split("\n"):
                    line = line.strip()
                    if not line:
                        continue

                    if line.startswith("■"):
                        # 목차(■) - 첫 번째가 아니면 앞에 빈 줄 추가
                        self._add_content_paragraph(
                            doc, line, is_highlight=True, add_space_before=not is_first_heading
                        )
                        is_first_heading = False
                    elif line.startswith("•") or line.startswith("-"):
                        # 불릿 포인트
                        clean_line = line.lstrip("•- ")
                        self._add_content_paragraph(doc, clean_line, is_bullet=True)
                    else:
                        self._add_content_paragraph(doc, line)

            # 차트 (있는 경우)
            if chart_key and chart_paths and chart_key in chart_paths:
                doc.add_paragraph()
                self._add_chart_image(doc, chart_paths[chart_key])

            # 페이지 브레이크 (마지막 섹션 제외)
            if section_id < len(sections):
                doc.add_page_break()

        # 4. 참고자료 (있는 경우)
        references = report_data.get("references")
        if references:
            self._add_section_heading(doc, len(sections) + 1, "참고자료 (References)")
            for line in references.split("\n"):
                if line.strip():
                    self._add_content_paragraph(doc, line.strip())

        # 5. 푸터
        self._add_footer(doc)

        # 저장
        if not output_filename:
            output_filename = f"AMORE_Analyst_Report_{start_date}_{end_date}.docx"

        output_path = self.output_dir / output_filename
        doc.save(output_path)

        logger.info(f"DOCX report generated: {output_path}")
        return output_path

    def to_bytes(self, doc_path: Path = None, doc: Document = None) -> bytes:
        """
        문서를 바이트로 변환

        Args:
            doc_path: 파일 경로 (doc이 None인 경우)
            doc: Document 객체 (doc_path가 None인 경우)

        Returns:
            문서 바이트
        """
        if doc_path:
            with open(doc_path, "rb") as f:
                return f.read()
        elif doc:
            buffer = BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.read()
        else:
            raise ValueError("Either doc_path or doc must be provided")


# =============================================================================
# PPTX Report Generator
# =============================================================================


class PptxReportGenerator:
    """
    PPTX 리포트 생성기

    AMOREPACIFIC IR 스타일의 프레젠테이션 생성
    """

    def __init__(self, output_dir: str = None):
        self.output_dir = Path(output_dir) if output_dir else Path(tempfile.gettempdir())
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # python-pptx 가용성 확인
        try:
            import importlib.util

            self._pptx_available = importlib.util.find_spec("pptx") is not None
        except Exception:
            self._pptx_available = False

        if not self._pptx_available:
            logger.warning("python-pptx not installed. PPTX generation disabled.")

        # Project root 찾기
        self.project_root = self._find_project_root()

    def _find_project_root(self) -> Path:
        """프로젝트 루트 디렉토리 찾기"""
        current = Path(__file__).resolve()
        for parent in current.parents:
            if (parent / "CLAUDE.md").exists() or (parent / "dashboard_api.py").exists():
                return parent
        return Path.cwd()

    def _get_logo_path(self, logo_type: str = "reverse") -> Path | None:
        """로고 파일 경로 반환 (PPTX는 어두운 배경이므로 reverse 권장)"""
        logo_map = {
            "color": DesignSystem.LOGO_COLOR_EN,
            "basic": DesignSystem.LOGO_BASIC_EN,
            "reverse": DesignSystem.LOGO_REVERSE_EN,
        }
        logo_rel = logo_map.get(logo_type, DesignSystem.LOGO_REVERSE_EN)
        logo_path = self.project_root / logo_rel

        if logo_path.exists():
            return logo_path
        return None

    def generate_presentation(
        self,
        report_data: dict[str, Any],
        chart_paths: dict[str, Path] = None,
        output_filename: str = None,
    ) -> Path | None:
        """
        프레젠테이션 생성

        Args:
            report_data: 리포트 데이터 (DocxReportGenerator와 동일 형식)
            chart_paths: 차트 이미지 경로
            output_filename: 출력 파일명

        Returns:
            생성된 파일 경로 (python-pptx 미설치 시 None)
        """
        if not self._pptx_available:
            logger.error("python-pptx not installed")
            return None

        from pptx import Presentation
        from pptx.dml.color import RGBColor as PptxRGBColor
        from pptx.util import Inches as PptxInches
        from pptx.util import Pt as PptxPt

        prs = Presentation()
        prs.slide_width = PptxInches(13.333)  # 16:9 비율
        prs.slide_height = PptxInches(7.5)

        # 색상 상수
        PACIFIC_BLUE = PptxRGBColor(0, 28, 88)
        WHITE = PptxRGBColor(255, 255, 255)
        GRAY = PptxRGBColor(125, 125, 125)
        _ = PACIFIC_BLUE, WHITE, GRAY  # 사용 표시

        # 메타 정보
        title = report_data.get("title", "AMORE Insight Report")
        start_date = report_data.get("start_date", "")
        end_date = report_data.get("end_date", "")

        # ========== 슬라이드 1: 표지 ==========
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # 배경 (Pacific Blue)
        background = slide.shapes.add_shape(
            1,  # Rectangle
            PptxInches(0),
            PptxInches(0),
            prs.slide_width,
            prs.slide_height,
        )
        background.fill.solid()
        background.fill.fore_color.rgb = PACIFIC_BLUE
        background.line.fill.background()

        # 로고
        logo_path = self._get_logo_path("reverse")
        if logo_path:
            slide.shapes.add_picture(
                str(logo_path), PptxInches(0.5), PptxInches(0.5), width=PptxInches(2.5)
            )

        # 제목
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(3), PptxInches(12), PptxInches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PptxPt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 날짜
        date_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(5), PptxInches(6), PptxInches(0.5)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{start_date} ~ {end_date}" if start_date else datetime.now().strftime("%Y-%m-%d")
        p.font.size = PptxPt(18)
        p.font.color.rgb = WHITE

        # ========== 슬라이드 2: 목차 ==========
        slide = prs.slides.add_slide(slide_layout)

        # AMORE PACIFIC 로고 이미지
        logo_path = self._get_logo_path("color")
        if logo_path:
            slide.shapes.add_picture(
                str(logo_path), PptxInches(0.5), PptxInches(0.3), width=PptxInches(3.5)
            )

        # 목차 항목 (IR 스타일 - 왼쪽 파란 바)
        sections = report_data.get("sections", [])
        y_pos = 1.5  # 시작 Y 위치 (로고 아래)

        for _i, section in enumerate(sections):
            # 왼쪽 파란 바 (Rectangle)
            bar = slide.shapes.add_shape(
                1,  # Rectangle
                PptxInches(0.8),
                PptxInches(y_pos),
                PptxInches(0.08),
                PptxInches(0.4),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = PACIFIC_BLUE
            bar.line.fill.background()

            # 섹션명 텍스트
            section_box = slide.shapes.add_textbox(
                PptxInches(1.0), PptxInches(y_pos), PptxInches(10), PptxInches(0.5)
            )
            tf = section_box.text_frame
            p = tf.paragraphs[0]
            p.text = section.get("title", "")
            p.font.size = PptxPt(18)
            p.font.bold = True
            p.font.color.rgb = PACIFIC_BLUE

            y_pos += 0.6  # 다음 항목 위치

        # DISCLAIMER (하단)
        disclaimer_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(6.5), PptxInches(12), PptxInches(0.8)
        )
        tf = disclaimer_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "DISCLAIMER"
        p.font.size = PptxPt(10)
        p.font.bold = True
        p.font.color.rgb = GRAY

        p2 = tf.add_paragraph()
        p2.text = (
            "본 자료는 Amazon.com 웹사이트에서 공개적으로 수집된 베스트셀러 순위 데이터를 기반으로 작성되었습니다. "
            "모든 분석 인사이트는 AI(인공지능)에 의해 자동 생성된 것으로, 실제 시장 상황과 다를 수 있습니다. "
            "본 자료의 정보를 활용한 의사결정에 대해 당사는 어떠한 책임도 지지 않습니다."
        )
        p2.font.size = PptxPt(8)
        p2.font.color.rgb = GRAY

        # ========== 본문 슬라이드들 ==========
        for section in sections:
            slide = prs.slides.add_slide(slide_layout)

            # 섹션 제목
            title_box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{section.get('id', '')}. {section.get('title', '')}"
            p.font.size = PptxPt(28)
            p.font.bold = True
            p.font.color.rgb = PACIFIC_BLUE

            # 내용 또는 차트
            chart_key = section.get("chart_key")
            if chart_key and chart_paths and chart_key in chart_paths:
                chart_path = chart_paths[chart_key]
                if Path(chart_path).exists():
                    slide.shapes.add_picture(
                        str(chart_path), PptxInches(1), PptxInches(1.5), width=PptxInches(11)
                    )
            else:
                # 텍스트 내용
                content = section.get("content", "")
                if content:
                    content_box = slide.shapes.add_textbox(
                        PptxInches(0.5), PptxInches(1.5), PptxInches(12), PptxInches(5.5)
                    )
                    tf = content_box.text_frame
                    tf.word_wrap = True

                    lines = content.split("\n")[:15]  # 최대 15줄
                    for i, line in enumerate(lines):
                        line = line.strip()
                        if not line:
                            continue
                        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                        p.text = line
                        p.font.size = PptxPt(16)
                        p.font.color.rgb = PACIFIC_BLUE if line.startswith("■") else GRAY
                        p.space_before = PptxPt(8)

        # 저장
        if not output_filename:
            output_filename = f"AMORE_Presentation_{start_date}_{end_date}.pptx"

        output_path = self.output_dir / output_filename
        prs.save(output_path)

        logger.info(f"PPTX presentation generated: {output_path}")
        return output_path


# =============================================================================
# PDF Report Generator
# =============================================================================


class PdfReportGenerator:
    """
    PDF 리포트 생성기

    DOCX를 PDF로 변환하거나 WeasyPrint로 직접 생성
    """

    def __init__(self, output_dir: str = None):
        self.output_dir = Path(output_dir) if output_dir else Path(tempfile.gettempdir())
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # WeasyPrint 가용성 확인
        try:
            import importlib.util

            self._weasyprint_available = importlib.util.find_spec("weasyprint") is not None
        except Exception:
            self._weasyprint_available = False

        if not self._weasyprint_available:
            logger.warning("WeasyPrint not installed. PDF generation limited.")

    def generate_from_html(
        self,
        html_content: str,
        output_filename: str = None,
    ) -> Path | None:
        """
        HTML에서 PDF 생성

        Args:
            html_content: HTML 문자열
            output_filename: 출력 파일명

        Returns:
            생성된 파일 경로
        """
        if not self._weasyprint_available:
            logger.error("WeasyPrint not installed")
            return None

        from weasyprint import CSS, HTML

        if not output_filename:
            output_filename = f"AMORE_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"

        output_path = self.output_dir / output_filename

        # CSS 스타일
        css = CSS(
            string="""
            @page {
                size: A4;
                margin: 2cm;
            }
            body {
                font-family: Arial, 'Noto Sans KR', sans-serif;
                font-size: 11pt;
                line-height: 1.5;
                color: #333333;
            }
            h1 {
                color: #001C58;
                font-size: 24pt;
                border-bottom: 2px solid #001C58;
                padding-bottom: 8px;
            }
            h2 {
                color: #1F5795;
                font-size: 16pt;
            }
            .highlight {
                color: #001C58;
                font-weight: bold;
            }
            .kpi-card {
                background: #F5F5F5;
                padding: 16px;
                border-radius: 8px;
                text-align: center;
            }
            .kpi-value {
                font-size: 28pt;
                font-weight: bold;
                color: #001C58;
            }
            .footer {
                text-align: center;
                font-size: 9pt;
                color: #7D7D7D;
                margin-top: 32px;
            }
        """
        )

        HTML(string=html_content).write_pdf(output_path, stylesheets=[css])

        logger.info(f"PDF generated: {output_path}")
        return output_path

    def convert_docx_to_pdf(self, docx_path: Path) -> Path | None:
        """
        DOCX를 PDF로 변환 (LibreOffice 또는 docx2pdf 사용)

        Args:
            docx_path: DOCX 파일 경로

        Returns:
            생성된 PDF 경로 (변환 실패 시 None)
        """
        import subprocess

        output_path = docx_path.with_suffix(".pdf")

        # 1. docx2pdf 시도 (Windows/macOS)
        try:
            from docx2pdf import convert

            convert(str(docx_path), str(output_path))
            if output_path.exists():
                logger.info(f"PDF converted via docx2pdf: {output_path}")
                return output_path
        except ImportError:
            pass
        except Exception as e:
            logger.warning(f"docx2pdf failed: {e}")

        # 2. LibreOffice 시도 (Linux/Docker)
        try:
            subprocess.run(
                [
                    "libreoffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(self.output_dir),
                    str(docx_path),
                ],
                check=True,
                capture_output=True,
            )
            if output_path.exists():
                logger.info(f"PDF converted via LibreOffice: {output_path}")
                return output_path
        except (subprocess.CalledProcessError, FileNotFoundError) as e:
            logger.warning(f"LibreOffice conversion failed: {e}")

        logger.error("PDF conversion failed: no converter available")
        return None


# =============================================================================
# Unified Report Generator
# =============================================================================


class ReportGenerator:
    """
    통합 리포트 생성기

    DOCX, PPTX, PDF 포맷을 모두 지원하는 통합 인터페이스
    """

    def __init__(self, output_dir: str = None):
        self.output_dir = Path(output_dir) if output_dir else Path(tempfile.gettempdir())
        self.output_dir.mkdir(parents=True, exist_ok=True)

        self.docx_generator = DocxReportGenerator(output_dir)
        self.pptx_generator = PptxReportGenerator(output_dir)
        self.pdf_generator = PdfReportGenerator(output_dir)

    def generate(
        self,
        report_data: dict[str, Any],
        chart_paths: dict[str, Path] = None,
        formats: list[str] = None,
        base_filename: str = None,
    ) -> dict[str, Path]:
        """
        지정된 포맷으로 리포트 생성

        Args:
            report_data: 리포트 데이터
            chart_paths: 차트 이미지 경로
            formats: 생성할 포맷 목록 ["docx", "pptx", "pdf"]
            base_filename: 기본 파일명 (확장자 제외)

        Returns:
            {"docx": Path, "pptx": Path, "pdf": Path}
        """
        if formats is None:
            formats = ["docx"]

        start_date = report_data.get("start_date", "")
        end_date = report_data.get("end_date", "")

        if not base_filename:
            base_filename = f"AMORE_Report_{start_date}_{end_date}"

        results = {}

        # DOCX 생성
        if "docx" in formats:
            docx_path = self.docx_generator.generate_analyst_report(
                report_data,
                chart_paths,
                output_filename=f"{base_filename}.docx",
            )
            results["docx"] = docx_path

        # PPTX 생성
        if "pptx" in formats:
            pptx_path = self.pptx_generator.generate_presentation(
                report_data,
                chart_paths,
                output_filename=f"{base_filename}.pptx",
            )
            if pptx_path:
                results["pptx"] = pptx_path

        # PDF 생성 (DOCX 변환)
        if "pdf" in formats:
            if "docx" in results:
                pdf_path = self.pdf_generator.convert_docx_to_pdf(results["docx"])
                if pdf_path:
                    results["pdf"] = pdf_path
            else:
                # DOCX 없이 PDF만 요청된 경우, 먼저 DOCX 생성 후 변환
                temp_docx = self.docx_generator.generate_analyst_report(
                    report_data,
                    chart_paths,
                    output_filename=f"{base_filename}_temp.docx",
                )
                pdf_path = self.pdf_generator.convert_docx_to_pdf(temp_docx)
                if pdf_path:
                    results["pdf"] = pdf_path
                # 임시 DOCX 삭제
                temp_docx.unlink(missing_ok=True)

        return results
